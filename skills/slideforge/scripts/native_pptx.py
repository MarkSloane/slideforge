"""
SlideForge Native PPTX Toolkit
===============================
Reusable helpers for generating editable PowerPoint slides with:
- Brand theme loading from ~/.slideforge/themes/
- Gradient backgrounds, cards, bars, bullet lists
- Fade entrance animations with staggered timing
- Apple-style drop shadows (boxes only, ~90% transparency)
- No shadows on lines or text (enforced)
- Horizontal gradient lines and bar fills

Usage in a presentation script:
    from native_pptx import SlideForgeNative
    sf = SlideForgeNative()  # auto-loads theme from ~/.slideforge/themes/
    slide = sf.add_slide()
    sf.add_textbox(slide, ...)
    sf.add_card(slide, ...)
    sf.add_gradient_line(slide, ...)
    sf.add_fade(slide, shape, delay_ms=100, duration_ms=125)
    sf.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import json
import glob
import os


class SlideForgeNative:
    """Native PPTX generator with brand theming, animations, and shadow control."""

    # ── Defaults (overridden by theme file) ──
    DEFAULTS = {
        "font": "Inter",
        "font_weights": "300;400;500;600;700",
        "primary_color": "#1a1a2e",
        "accent_color": "#e94560",
        "accent_gradient": "linear-gradient(90deg, #e94560, #c23152)",
        "secondary_text": "#6b7280",
        "background": "radial-gradient(ellipse at center, #f0f0f0 0%, #d8d8d8 100%)",
    }

    def __init__(self, theme_path=None):
        """Initialize with optional theme. Auto-loads from ~/.slideforge/themes/ if not specified."""
        self.theme = dict(self.DEFAULTS)
        self._load_theme(theme_path)
        self._parse_colors()

        # Presentation object
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # Layout constants
        self.SLIDE_W = self.prs.slide_width
        self.SLIDE_H = self.prs.slide_height
        self.M_LEFT = Inches(1.0)
        self.M_RIGHT = Inches(1.0)
        self.CONTENT_W = self.SLIDE_W - self.M_LEFT - self.M_RIGHT

    def _load_theme(self, theme_path):
        """Load theme from file or auto-detect from ~/.slideforge/themes/."""
        if theme_path and os.path.exists(theme_path):
            with open(theme_path) as f:
                self.theme.update(json.load(f))
            return

        themes_dir = os.path.expanduser("~/.slideforge/themes")
        if os.path.isdir(themes_dir):
            files = sorted(glob.glob(os.path.join(themes_dir, "*.json")))
            if files:
                with open(files[0]) as f:
                    self.theme.update(json.load(f))

    def _hex_to_rgb(self, hex_str):
        """Convert '#RRGGBB' to RGBColor."""
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _parse_gradient_colors(self, grad_str):
        """Extract two hex colors from 'linear-gradient(90deg, #aaa, #bbb)'."""
        import re
        colors = re.findall(r'#[0-9A-Fa-f]{6}', grad_str)
        if len(colors) >= 2:
            return self._hex_to_rgb(colors[0]), self._hex_to_rgb(colors[1])
        c = self._hex_to_rgb(colors[0]) if colors else self.accent
        return c, c

    def _parse_bg_colors(self, bg_str):
        """Extract two hex colors from background gradient string."""
        import re
        colors = re.findall(r'#[0-9A-Fa-f]{6}', bg_str)
        if len(colors) >= 2:
            return self._hex_to_rgb(colors[0]), self._hex_to_rgb(colors[1])
        c = self._hex_to_rgb(colors[0]) if colors else RGBColor(0xF0, 0xF0, 0xF0)
        return c, c

    def _parse_colors(self):
        """Parse theme into RGBColor attributes."""
        self.primary = self._hex_to_rgb(self.theme["primary_color"])
        self.accent = self._hex_to_rgb(self.theme["accent_color"])
        self.secondary_text = self._hex_to_rgb(self.theme["secondary_text"])
        self.accent_grad_start, self.accent_grad_end = self._parse_gradient_colors(
            self.theme["accent_gradient"])
        self.bg_start, self.bg_end = self._parse_bg_colors(self.theme["background"])
        self.font = self.theme["font"]

        # Derived colors
        self.white = RGBColor(0xFF, 0xFF, 0xFF)
        self.light_gray = RGBColor(0x9C, 0xA3, 0xAF)

        # Extra palette colors (if present in theme)
        palette = self.theme.get("palette", {})
        self.palette = {k: self._hex_to_rgb(v) for k, v in palette.items()}

    def color(self, name):
        """Get a palette color by name, falling back to accent."""
        return self.palette.get(name, self.accent)

    # ═══════════════════════════════════════════════
    # Slide & Background
    # ═══════════════════════════════════════════════

    def add_slide(self):
        """Add a blank slide with gradient background."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = self.bg_start
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = self.bg_end
        fill.gradient_stops[1].position = 1.0
        return slide

    # ═══════════════════════════════════════════════
    # Text
    # ═══════════════════════════════════════════════

    def add_textbox(self, slide, left, top, width, height, text,
                    font_size=Pt(14), color=None, bold=False, italic=False,
                    alignment=PP_ALIGN.LEFT):
        """Add a simple text box. Returns (shape, paragraph)."""
        if color is None:
            color = self.primary
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.font
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = alignment
        return txBox, p

    def add_multi_text(self, slide, left, top, width, height, runs,
                       alignment=PP_ALIGN.LEFT):
        """Add a textbox with multiple styled runs.
        runs = [(text, font_size, color, bold, italic), ...]
        """
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = alignment
        for i, (text, font_size, color, bold, italic) in enumerate(runs):
            if i == 0:
                run = p.runs[0] if p.runs else p.add_run()
                run.text = text
            else:
                run = p.add_run()
                run.text = text
            run.font.name = self.font
            run.font.size = font_size
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.italic = italic
        return txBox

    def add_bullet_list(self, slide, left, top, width, height, items,
                        font_size=Pt(13), color=None, bullet_color=None,
                        line_spacing=Pt(24)):
        """Add a bulleted list text box."""
        if color is None:
            color = self.secondary_text
        if bullet_color is None:
            bullet_color = self.accent
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = line_spacing
            run_bullet = p.add_run()
            run_bullet.text = "\u25cf  "
            run_bullet.font.name = self.font
            run_bullet.font.size = Pt(8)
            run_bullet.font.color.rgb = bullet_color
            run_text = p.add_run()
            run_text.text = item
            run_text.font.name = self.font
            run_text.font.size = font_size
            run_text.font.color.rgb = color
        return txBox

    # ═══════════════════════════════════════════════
    # Shapes
    # ═══════════════════════════════════════════════

    def add_card(self, slide, left, top, width, height, corner_ratio=0.04):
        """Add a rounded rectangle card with subtle Apple-style shadow.
        Shadow: boxes only, 90% transparent, 3pt blur, 2pt offset down.
        """
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.white
        shape.line.fill.background()
        shape.adjustments[0] = corner_ratio
        self._add_subtle_shadow(shape)
        return shape

    def add_gradient_line(self, slide, left, top, width, height=Pt(3)):
        """Add accent gradient line (thin rectangle). No shadow ever."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height)
        fill = shape.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = self.accent_grad_start
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = self.accent_grad_end
        fill.gradient_stops[1].position = 1.0
        shape.line.fill.background()
        shape.rotation = 0.0
        self._set_gradient_horizontal(shape)
        self._remove_shadow(shape)
        return shape

    def add_gradient_bar(self, slide, left, top, width, height,
                         corner_ratio=0.3):
        """Add a gradient-filled bar (for charts). No shadow."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        fill = shape.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = self.accent_grad_start
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = self.accent_grad_end
        fill.gradient_stops[1].position = 1.0
        shape.line.fill.background()
        shape.adjustments[0] = corner_ratio
        self._set_gradient_horizontal(shape)
        self._remove_shadow(shape)
        return shape

    def add_solid_bar(self, slide, left, top, width, height,
                      color=None, corner_ratio=0.3):
        """Add a solid-color bar (for charts). No shadow."""
        if color is None:
            color = RGBColor(0xB0, 0xB8, 0xC4)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.adjustments[0] = corner_ratio
        self._remove_shadow(shape)
        return shape

    def add_bar_track(self, slide, left, top, width, height, corner_ratio=0.3):
        """Add a bar chart track (background). No shadow."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.bg_end
        shape.line.fill.background()
        shape.adjustments[0] = corner_ratio
        self._remove_shadow(shape)
        return shape

    def add_badge(self, slide, left, top, width, height, text,
                  bg_color=None, text_color=None, corner_ratio=0.3):
        """Add a small colored badge with centered text. No shadow."""
        if bg_color is None:
            # Light tint of accent
            r = min(255, self.accent_grad_start[0] + 200)
            g = min(255, self.accent_grad_start[1] + 200)
            b = min(255, self.accent_grad_start[2] + 200)
            bg_color = RGBColor(r, g, b)
        if text_color is None:
            text_color = self.accent
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        shape.adjustments[0] = corner_ratio
        self._remove_shadow(shape)
        self.add_textbox(slide, left, top, width, height, text,
                         font_size=Pt(9), color=text_color, bold=True,
                         alignment=PP_ALIGN.CENTER)
        return shape

    # ═══════════════════════════════════════════════
    # Shadow Control
    # ═══════════════════════════════════════════════

    @staticmethod
    def _add_subtle_shadow(shape):
        """Apple-style drop shadow: 3pt blur, 2pt offset, 90% transparent.
        ONLY for card/box shapes. Never call on lines or text.
        """
        spPr = shape._element.spPr
        for child in spPr.findall(qn('a:effectLst')):
            spPr.remove(child)
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', '38100')    # 3pt
        outerShdw.set('dist', '25400')        # 2pt
        outerShdw.set('dir', '5400000')       # straight down
        outerShdw.set('algn', 'tl')
        outerShdw.set('rotWithShape', '0')
        srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '10000')  # 10% opacity

    @staticmethod
    def _remove_shadow(shape):
        """Explicitly strip shadow and prevent PowerPoint defaults."""
        spPr = shape._element.spPr
        for child in spPr.findall(qn('a:effectLst')):
            spPr.remove(child)
        etree.SubElement(spPr, qn('a:effectLst'))  # empty = no effects

    # ═══════════════════════════════════════════════
    # Gradient Direction
    # ═══════════════════════════════════════════════

    @staticmethod
    def _set_gradient_horizontal(shape):
        """Force gradient left-to-right (0 degrees)."""
        gradFill = shape._element.spPr.find(qn('a:gradFill'))
        if gradFill is None:
            return
        for lin in gradFill.findall(qn('a:lin')):
            gradFill.remove(lin)
        lin = etree.SubElement(gradFill, qn('a:lin'))
        lin.set('ang', '0')
        lin.set('scaled', '1')

    # ═══════════════════════════════════════════════
    # Animations
    # ═══════════════════════════════════════════════

    @staticmethod
    def _get_or_create_timing(slide):
        """Get or create the animation timing tree for a slide."""
        sld = slide._element
        timing = sld.find(qn('p:timing'))
        if timing is None:
            timing = etree.SubElement(sld, qn('p:timing'))
        tnLst = timing.find(qn('p:tnLst'))
        if tnLst is None:
            tnLst = etree.SubElement(timing, qn('p:tnLst'))
        par = tnLst.find(qn('p:par'))
        if par is None:
            par = etree.SubElement(tnLst, qn('p:par'))
            cTn = etree.SubElement(par, qn('p:cTn'))
            cTn.set('id', '1')
            cTn.set('dur', 'indefinite')
            cTn.set('restart', 'never')
            cTn.set('nodeType', 'tmRoot')
        return par

    @staticmethod
    def _next_id(slide):
        """Get next unique timing node ID."""
        sld = slide._element
        ids = [int(e.get('id', '0')) for e in sld.iter()
               if e.get('id') is not None]
        return max(ids) + 1 if ids else 2

    def add_fade(self, slide, shape, delay_ms=0, duration_ms=125):
        """Add a fade entrance animation to a shape.
        Default timing is snappy (125ms duration).
        """
        par_root = self._get_or_create_timing(slide)
        cTn_root = par_root.find(qn('p:cTn'))
        childTnLst = cTn_root.find(qn('p:childTnLst'))
        if childTnLst is None:
            childTnLst = etree.SubElement(cTn_root, qn('p:childTnLst'))

        seq = childTnLst.find(qn('p:seq'))
        if seq is None:
            seq = etree.SubElement(childTnLst, qn('p:seq'))
            seq.set('concurrent', '1')
            seq.set('nextAc', 'seek')
            seq_cTn = etree.SubElement(seq, qn('p:cTn'))
            nid = self._next_id(slide)
            seq_cTn.set('id', str(nid))
            seq_cTn.set('dur', 'indefinite')
            seq_cTn.set('nodeType', 'mainSeq')
            seq_childTnLst = etree.SubElement(seq_cTn, qn('p:childTnLst'))

            prevCondLst = etree.SubElement(seq, qn('p:prevCondLst'))
            cond_prev = etree.SubElement(prevCondLst, qn('p:cond'))
            cond_prev.set('evt', 'onPrev')
            cond_prev.set('delay', '0')
            tgtEl_prev = etree.SubElement(cond_prev, qn('p:tgtEl'))
            etree.SubElement(tgtEl_prev, qn('p:sldTgt'))

            nextCondLst = etree.SubElement(seq, qn('p:nextCondLst'))
            cond_next = etree.SubElement(nextCondLst, qn('p:cond'))
            cond_next.set('evt', 'onNext')
            cond_next.set('delay', '0')
            tgtEl_next = etree.SubElement(cond_next, qn('p:tgtEl'))
            etree.SubElement(tgtEl_next, qn('p:sldTgt'))
        else:
            seq_cTn = seq.find(qn('p:cTn'))
            seq_childTnLst = seq_cTn.find(qn('p:childTnLst'))
            if seq_childTnLst is None:
                seq_childTnLst = etree.SubElement(seq_cTn, qn('p:childTnLst'))

        # Animation group
        par1 = etree.SubElement(seq_childTnLst, qn('p:par'))
        cTn1 = etree.SubElement(par1, qn('p:cTn'))
        nid = self._next_id(slide)
        cTn1.set('id', str(nid))
        cTn1.set('fill', 'hold')
        stCondLst1 = etree.SubElement(cTn1, qn('p:stCondLst'))
        cond1 = etree.SubElement(stCondLst1, qn('p:cond'))
        cond1.set('delay', '0')

        childTnLst1 = etree.SubElement(cTn1, qn('p:childTnLst'))
        par2 = etree.SubElement(childTnLst1, qn('p:par'))
        cTn2 = etree.SubElement(par2, qn('p:cTn'))
        nid = self._next_id(slide)
        cTn2.set('id', str(nid))
        cTn2.set('fill', 'hold')
        stCondLst2 = etree.SubElement(cTn2, qn('p:stCondLst'))
        cond2 = etree.SubElement(stCondLst2, qn('p:cond'))
        cond2.set('delay', str(delay_ms))

        childTnLst2 = etree.SubElement(cTn2, qn('p:childTnLst'))
        par3 = etree.SubElement(childTnLst2, qn('p:par'))
        cTn3 = etree.SubElement(par3, qn('p:cTn'))
        nid = self._next_id(slide)
        cTn3.set('id', str(nid))
        cTn3.set('presetID', '10')  # Fade
        cTn3.set('presetClass', 'entr')
        cTn3.set('presetSubtype', '0')
        cTn3.set('fill', 'hold')
        cTn3.set('grpId', '0')
        cTn3.set('nodeType', 'withEffect')
        stCondLst3 = etree.SubElement(cTn3, qn('p:stCondLst'))
        cond3 = etree.SubElement(stCondLst3, qn('p:cond'))
        cond3.set('delay', '0')
        childTnLst3 = etree.SubElement(cTn3, qn('p:childTnLst'))

        sp_id = shape.shape_id

        # Fade effect
        animEffect = etree.SubElement(childTnLst3, qn('p:animEffect'))
        animEffect.set('transition', 'in')
        animEffect.set('filter', 'fade')
        ae_cBhvr = etree.SubElement(animEffect, qn('p:cBhvr'))
        ae_cTn = etree.SubElement(ae_cBhvr, qn('p:cTn'))
        nid = self._next_id(slide)
        ae_cTn.set('id', str(nid))
        ae_cTn.set('dur', str(duration_ms))
        ae_tgtEl = etree.SubElement(ae_cBhvr, qn('p:tgtEl'))
        ae_spTgt = etree.SubElement(ae_tgtEl, qn('p:spTgt'))
        ae_spTgt.set('spid', str(sp_id))

        # Visibility set
        st = etree.SubElement(childTnLst3, qn('p:set'))
        st_cBhvr = etree.SubElement(st, qn('p:cBhvr'))
        st_cTn = etree.SubElement(st_cBhvr, qn('p:cTn'))
        nid = self._next_id(slide)
        st_cTn.set('id', str(nid))
        st_cTn.set('dur', '1')
        st_cTn.set('fill', 'hold')
        stCondLst_s = etree.SubElement(st_cTn, qn('p:stCondLst'))
        cond_s = etree.SubElement(stCondLst_s, qn('p:cond'))
        cond_s.set('delay', '0')
        st_tgtEl = etree.SubElement(st_cBhvr, qn('p:tgtEl'))
        st_spTgt = etree.SubElement(st_tgtEl, qn('p:spTgt'))
        st_spTgt.set('spid', str(sp_id))
        st_attrNameLst = etree.SubElement(st_cBhvr, qn('p:attrNameLst'))
        st_attrName = etree.SubElement(st_attrNameLst, qn('p:attrName'))
        st_attrName.text = 'style.visibility'
        st_to = etree.SubElement(st, qn('p:to'))
        st_val = etree.SubElement(st_to, qn('p:strVal'))
        st_val.set('val', 'visible')

    def add_fade_sequence(self, slide, shapes_and_delays):
        """Convenience: animate a list of (shape, delay_ms) pairs.
        Example: sf.add_fade_sequence(slide, [(heading, 0), (divider, 50), (card, 100)])
        """
        for shape, delay in shapes_and_delays:
            self.add_fade(slide, shape, delay_ms=delay)

    def add_fade_stagger(self, slide, shapes, start_delay=100, gap=60,
                         duration_ms=125):
        """Convenience: stagger-animate a list of shapes.
        Example: sf.add_fade_stagger(slide, [card1, card2, card3])
        """
        for i, shape in enumerate(shapes):
            self.add_fade(slide, shape,
                          delay_ms=start_delay + i * gap,
                          duration_ms=duration_ms)

    # ═══════════════════════════════════════════════
    # Save
    # ═══════════════════════════════════════════════

    def save(self, path):
        """Save the presentation."""
        self.prs.save(path)
        print(f"Saved: {path}")
