#!/usr/bin/env python3
"""
SlideForge PPTX Exporter
Screenshots HTML slides via Playwright and assembles a pixel-perfect PPTX.

Usage:
    python3 html_to_pptx.py <slides_dir> <output_pptx> [--title "Presentation Title"]

Requires: playwright, python-pptx
    pip3 install playwright python-pptx
    python3 -m playwright install chromium
"""

import argparse
import sys
import time
from pathlib import Path

def main():
    parser = argparse.ArgumentParser(description="Convert SlideForge HTML slides to PPTX")
    parser.add_argument("slides_dir", help="Directory containing slide-01.html, slide-02.html, etc.")
    parser.add_argument("output", help="Output .pptx file path")
    parser.add_argument("--title", default="Presentation", help="Presentation title (metadata)")
    parser.add_argument("--wait", type=float, default=5.0,
                        help="Seconds to wait for animations before capture (default: 5)")
    parser.add_argument("--width", type=int, default=1920, help="Viewport width (default: 1920)")
    parser.add_argument("--height", type=int, default=1080, help="Viewport height (default: 1080)")
    parser.add_argument("--scale", type=int, default=2, help="Device scale factor for retina (default: 2)")
    args = parser.parse_args()

    slides_dir = Path(args.slides_dir).resolve()
    output_path = Path(args.output).resolve()

    # Discover slides
    slide_files = sorted(slides_dir.glob("slide-*.html"))
    if not slide_files:
        print(f"Error: No slide-*.html files found in {slides_dir}", file=sys.stderr)
        sys.exit(1)

    print(f"Found {len(slide_files)} slides in {slides_dir}")

    # Check dependencies
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print("Error: playwright not installed. Run: pip3 install playwright && python3 -m playwright install chromium",
              file=sys.stderr)
        sys.exit(1)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError:
        print("Error: python-pptx not installed. Run: pip3 install python-pptx", file=sys.stderr)
        sys.exit(1)

    # Capture screenshots
    temp_images = []
    print("Capturing slides...")

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={"width": args.width, "height": args.height},
            device_scale_factor=args.scale
        )

        for slide_file in slide_files:
            page.goto(f"file://{slide_file}")
            time.sleep(args.wait)
            img_path = slides_dir / f"_tmp_{slide_file.stem}.png"
            page.screenshot(path=str(img_path))
            temp_images.append(img_path)
            print(f"  Captured {slide_file.name}")

        browser.close()

    # Build PPTX
    print("Building PPTX...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for img_path in temp_images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        slide.shapes.add_picture(str(img_path), Emu(0), Emu(0), prs.slide_width, prs.slide_height)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    # Cleanup temp images
    for img_path in temp_images:
        img_path.unlink(missing_ok=True)

    print(f"Saved: {output_path}")
    print(f"Slides: {len(slide_files)}")

if __name__ == "__main__":
    main()
